"""
Auditoria geometrica del PPTX.

No hay LibreOffice en esta maquina, de modo que no se puede renderizar las
diapositivas a imagen. Este script comprueba lo que si es medible a partir del
XML: que ninguna forma salga de la diapositiva, que se respete el margen, que no
haya solapamientos de texto, y una estimacion de desborde vertical del texto a
partir de un ancho medio de caracter por tamano de fuente.

La estimacion de desborde es aproximada: avisa de casos gruesos, no sustituye a
mirar la diapositiva.
"""

import sys
from pptx import Presentation
from pptx.util import Emu

RUTA = sys.argv[1] if len(sys.argv) > 1 else "DEFENSA_TFM.pptx"
ANCHO, ALTO = 13.3, 7.5
MARGEN_MIN = 0.5

# Ancho medio de caracter como fraccion del tamano de fuente, por familia.
# Calibri y Cambria a cuerpo corrido rondan 0.48-0.52 em.
ANCHO_CAR = 0.50
ALTURA_LINEA = 1.22   # multiplicador tipico


def pulgadas(v):
    return Emu(v).inches if v is not None else None


def altura_estimada(forma, ancho_pulg, tamano_defecto=12):
    """Altura aproximada del texto, parrafo a parrafo.

    Usar el tamano mayor de la forma para TODAS las lineas sobreestima mucho
    cuando un parrafo mezcla una cifra grande con texto pequeno, que es el caso
    de las tarjetas de esta presentacion. Aqui cada parrafo aporta su propia
    altura de linea, tomada del run mas grande de ese parrafo, y el ancho medio
    de caracter se pondera por la longitud de cada run.
    """
    alto = 0.0
    detalle = []
    for par in forma.text_frame.paragraphs:
        runs = [(r.text or "", r.font.size.pt if r.font.size else tamano_defecto)
                for r in par.runs]
        if not runs:
            alto += tamano_defecto * ALTURA_LINEA / 72
            continue
        nchars = sum(len(t) for t, _ in runs) or 1
        pt_max = max(pt for _, pt in runs)
        # ancho medio de caracter ponderado por longitud de run
        ancho_medio = sum(len(t) * pt * ANCHO_CAR for t, pt in runs) / nchars
        car_por_linea = max(int((ancho_pulg * 72) / max(ancho_medio, 1)), 1)
        nl = max(1, -(-nchars // car_por_linea))
        alto += nl * pt_max * ALTURA_LINEA / 72
        detalle.append(f"{nl}L@{pt_max:.0f}pt")
    return alto, ", ".join(detalle)


pres = Presentation(RUTA)
problemas = []

for i, diapo in enumerate(pres.slides, start=1):
    cajas = []
    for f in diapo.shapes:
        x, y = pulgadas(f.left), pulgadas(f.top)
        w, h = pulgadas(f.width), pulgadas(f.height)
        if None in (x, y, w, h):
            continue

        nombre = f.shape_type
        txt = f.text_frame.text if f.has_text_frame else ""
        etq = (txt[:38].replace("\n", " ") + "…") if len(txt) > 38 else txt.replace("\n", " ")
        ident = f"{etq or nombre}"

        # 1. Fuera de la diapositiva
        if x < -0.01 or y < -0.01 or x + w > ANCHO + 0.01 or y + h > ALTO + 0.01:
            problemas.append(
                f"D{i:02d}  FUERA DE LIMITES  «{ident}»  "
                f"x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f} "
                f"(derecha {x + w:.2f}/{ANCHO}, abajo {y + h:.2f}/{ALTO})")

        # 2. Margen insuficiente
        elif x < MARGEN_MIN - 0.01 or y < 0.15 or ANCHO - (x + w) < MARGEN_MIN - 0.01:
            problemas.append(
                f"D{i:02d}  MARGEN ESCASO     «{ident}»  "
                f"izq={x:.2f} der={ANCHO - (x + w):.2f} arriba={y:.2f}")

        # 3. Desborde vertical estimado del texto
        if txt.strip() and f.has_text_frame:
            alto_nec, detalle = altura_estimada(f, w)
            if alto_nec > h * 1.12:
                problemas.append(
                    f"D{i:02d}  DESBORDE ¿?       «{ident}»  "
                    f"{detalle} necesitan ~{alto_nec:.2f}\" "
                    f"y la caja mide {h:.2f}\"")

        if txt.strip():
            cajas.append((ident, x, y, w, h))

    # 4. Solapamiento entre cajas de texto
    for a in range(len(cajas)):
        for b in range(a + 1, len(cajas)):
            n1, x1, y1, w1, h1 = cajas[a]
            n2, x2, y2, w2, h2 = cajas[b]
            sx = min(x1 + w1, x2 + w2) - max(x1, x2)
            sy = min(y1 + h1, y2 + h2) - max(y1, y2)
            if sx > 0.06 and sy > 0.06:
                problemas.append(
                    f"D{i:02d}  SOLAPAMIENTO      «{n1}» × «{n2}»  "
                    f"{sx:.2f}\" × {sy:.2f}\"")

print(f"Diapositivas: {len(pres.slides)}")
print(f"Incidencias:  {len(problemas)}\n")
for p in problemas:
    print(p)
if not problemas:
    print("Sin incidencias geometricas.")
